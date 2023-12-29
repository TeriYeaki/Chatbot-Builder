import PyPDF2
from docx import Document
from pptx import Presentation

def pdf_extract(pdf_path):
    """
    Extract text from pdf file using PyPDF2 library.
    """
    
    text = ""
    with open(pdf_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in range(len(pdf_reader.pages)):
            text += pdf_reader.pages[page].extract_text()
    return text


def docx_extract(file_path):
    """
    Extract text from docx file using docx library.
    """
    doc = Document(file_path)
    text = [paragraph.text for paragraph in doc.paragraphs]
    return '\n'.join(text)


def txt_extract(file_path):
    """
    Extract text from txt file.
    """
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return text


def pptx_extract(file_path):
    """
    Extract text from pptx file using pptx library.
    """
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)