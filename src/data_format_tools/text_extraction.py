from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer
from docx import Document
from pptx import Presentation
import os
import re

class Document_Text_Extraction():
    """
    Extract text from document files.
    """

    def __init__(self, file_path):
        """
        Initialize file path and file extension.
        """
        self.file_path = file_path
        self.file_extension = file_path.split(".")[-1]
        self.file_name = self._get_file_name()
    
    def _get_file_name(self):
        """
        Return file name.
        """
        base_name = os.path.basename(self.file_path)
        return os.path.splitext(base_name)[0]

    def _pdf_extract(self, header_footer_margin=50):
        """
        Extract raw text from pdf file using PDFMiner library. It ignores text in the header and footer.
        It extracts text from each page and joins them together.
        """

        print("Extracting text from pdf file...")

        text_content = []
        for page_layout in extract_pages(self.file_path):
            page_height = page_layout.height
            for element in page_layout:
                if isinstance(element, LTTextContainer):
                    # Check if the text element is outside the header/footer margin
                    if header_footer_margin < element.y0 < (page_height - header_footer_margin):
                        text_content.append(element.get_text())
        
        print("Text extraction complete.")

        return self._preprocess_text("\n".join(text_content))


    def _docx_extract(self):
            """
            Extracts text from a docx file using the docx library.

            Returns:
                str: The extracted text from the docx file.
            """

            print("Extracting text from docx file...")

            doc = Document(self.file_path)
            text = [paragraph.text for paragraph in doc.paragraphs]

            print("Text extraction complete.")

            return self._preprocess_text('\n'.join(text))


    def _txt_extract(self):
        """
        Extract text from txt file.
        """

        print("Extracting text from txt file...")

        with open(self.file_path, 'r', encoding='utf-8') as file:
            text = file.read()

        print("Text extraction complete.")

        return self._preprocess_text(text)


    def _pptx_extract(self):
        """
        Extract text from pptx file using pptx library.
        """

        print("Extracting text from pptx file...")

        prs = Presentation(self.file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

        print("Text extraction complete.")

        return self._preprocess_text(("\n".join(text)))

    def _preprocess_text(self, raw_text):
        """
        Basic preprocess text by removing unnecessary characters.
        """
        text = re.sub(r'\s+', ' ', raw_text)
        return text
    

    def extract_text(self):
        """
        Extract text from file based on file extension.
        """
        match self.file_extension:
            case 'pdf':
                print("File extension detected as pdf.")
                return self._pdf_extract()
            case'docx':
                print("File extension detected as docx.")
                return self._docx_extract()
            case'txt':
                print("File extension detected as txt.")
                return self._txt_extract()
            case 'pptx':
                print("File extension detected as pptx.")
                return self._pptx_extract()
            case _:
                raise Exception("Error: File extension not supported. Please use .pdf, .docx, .txt, or .pptx files.")
    
    